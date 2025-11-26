from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
slides_content = [
    ("Modernismo – 3ª Fase e Tropicalismo", "Análise literária e artística\nAluno(a): ______ | Data: ______"),
    ("Visão Geral", "Esta apresentação aborda a 3ª fase do Modernismo brasileiro e o movimento artístico Tropicalista."),
    ("Contexto Histórico", "• Início em 1945\n• Pós-guerra\n• Modernização e urbanização\n• Literatura reflexiva"),
    ("Características", "• Linguagem introspectiva\n• Temas existenciais\n• Inovação narrativa"),
    ("Autores", "Clarice Lispector\nGuimarães Rosa\nJoão Cabral de Melo Neto"),
    ("Clarice Lispector", "A Hora da Estrela (1977)"),
    ("Guimarães Rosa", "Grande Sertão: Veredas (1956)"),
    ("João Cabral de Melo Neto", "Morte e Vida Severina (1955)"),
    ("Tropicalismo", "Movimento de 1967–68, mistura cultural e crítica"),
    ("Características do Tropicalismo", "• Mistura de ritmos\n• Antropofagia\n• Crítica social"),
    ("Artistas", "Pintura: Oiticica, Lygia Clark\nPoesia: Torquato Neto\nMúsica: Caetano, Gil, Mutantes"),
    ("Conclusão", "O Modernismo 3ª fase e o Tropicalismo transformaram a arte brasileira."),
    ("Envio", "Enviar para: sidnei.felix@edu.se.df.gov.br")
]

for title, body in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body

fname = "apresentacao_modernismo_tropicalismo.pptx"
prs.save(fname)
print("Arquivo gerado:", fname)
